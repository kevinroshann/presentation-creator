from flask import Flask, request
import textwrap
from pptx import Presentation
from pptx.util import Pt  # Importing Pt class for font size
import google.generativeai as genai





import textwrap
from pptx import Presentation
from pptx.util import Pt  # Importing Pt class for font size
import google.generativeai as genai


def to_markdown(text):
  """Fixes indentation issue and converts text to markdown format with proper bullet points."""
  if isinstance(text, list):  # Check if content is a list
    text = '\n'.join(text)  # Join list elements into a single string
  # Use consistent indentation for bullets (two spaces)
  text = "\n".join([f"  * {line}" for line in text.split("\n")])
  return textwrap.indent(text, '> ', predicate=lambda _: True)

def generate_presentation(topic,num):
  GOOGLE_API_KEY = 'AIzaSyCsRvwDyPyi14PrVSq5w2tk_Uf2kqhGYMM'  # Replace with your actual Google API key
  genai.configure(api_key=GOOGLE_API_KEY)
  model = genai.GenerativeModel('gemini-pro')

  presentation = Presentation()
  layout = presentation.slide_layouts[1]  # Using the default layout for slides

  # Add title slide
  title_slide_layout = presentation.slide_layouts[0]
  title_slide = presentation.slides.add_slide(title_slide_layout)
  title = title_slide.shapes.title
  title.text = "Presentation on " + topic
  
  # Prompt user for confirmation on AI-generated subtopics
  n=int(num)
    # Generate subtopics using the main topic
  subtopic_prompt = f"Generate {n} subtopics of 5 words in each subtopic related to the topic {topic}"
  subtopic_response = model.generate_content(subtopic_prompt)
  subtopics = subtopic_response.text.split("\n")[:n]  # Get the first 5 lines
  
  # Get font sizes (optional, adjust as needed)
  font_sizes = [14] * len(subtopics)  # Set a default font size for all slides
 
  # Create slides for each subtopic
  for i, (subtopic, font_size) in enumerate(zip(subtopics, font_sizes)):
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = subtopic

      # Accessing the font object and setting size
    title_shape = slide.shapes.title
    title_font = title_shape.text_frame.paragraphs[0].runs[0].font
    title_font.size = Pt(32)
    
    

    # Generate content for the subtopic
    response = model.generate_content(f"generate a brief paragraph on the topic {subtopic}")
    content = response.text.split('\n')  # Split into lines

    # Add content to the slide
    textbox = slide.placeholders[1]
    text_frame = textbox.text_frame
    text_frame.text = to_markdown(content)

    # Set font size
    for paragraph in text_frame.paragraphs:
      for run in paragraph.runs:
        run.font.size = Pt(font_size)

  # Add conclusion slide
  conclusion_slide_layout = presentation.slide_layouts[1]  # Use the last slide layout
  conclusion_slide = presentation.slides.add_slide(conclusion_slide_layout)
  conclusion_slide.shapes.title.text = "Conclusion"
  response = model.generate_content(f"generate a 5 sentence conclusion summarizing {topic}")
  content = response.text.split('\n')  # Split into lines

  # Add content to the slide
  textbox = conclusion_slide.placeholders[1]
  text_frame = textbox.text_frame
  text_frame.text = to_markdown(content)
  for paragraph in text_frame.paragraphs:
      for run in paragraph.runs:
        run.font.size = Pt(font_size)

  # Save the presentation
  presentation.save(f"{topic}_presentation.pptx")

# Call the function to generate the presentation




app = Flask(__name__)

@app.route("/", methods=["GET"])
def index():
    # Serve the HTML file
    return open("index.html", "r").read()

@app.route("/process_data", methods=["POST"])
def process_data():
    # Get name and topic from request data
    name = request.form["name"]
    topic = request.form["topic"]
    num = request.form["num"]
    
    # Perform actions with the data (e.g., store in database, log to file)
    print(f"Name: {name}, Topic: {topic}")

    generate_presentation(topic,num)  # Example processing
    return "Data received successfully!", 200  # Return success message
    

if __name__ == "__main__":
    app.run(debug=True)  # Set debug=False for production
